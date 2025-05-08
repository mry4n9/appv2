import PyPDF2
from pptx import Presentation
import re

def add_http_if_missing(url):
    """Adds http:// or https:// to a URL if the scheme is missing."""
    if not url:
        return ""
    if not re.match(r'(?:http|ftp|https)://', url):
        return 'https://' + url  # Default to https
    return url

def extract_text_from_pdf(file_obj):
    """Extracts text from an uploaded PDF file object."""
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return ""

def extract_text_from_pptx(file_obj):
    """Extracts text from an uploaded PPTX file object."""
    try:
        prs = Presentation(file_obj)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return ""

def get_file_extension(filename):
    """Gets the file extension from a filename."""
    return filename.split('.')[-1].lower() if '.' in filename else ""

def format_company_name_for_filename(company_name):
    """Cleans a company name for use in filenames."""
    if not company_name:
        return "generic_company"
    # Remove special characters, replace spaces with underscores
    name = re.sub(r'[^\w\s-]', '', company_name).strip().replace(' ', '_')
    name = re.sub(r'[-\s]+', '_', name)
    return name if name else "generic_company"